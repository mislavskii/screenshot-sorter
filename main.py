import os
from pptx import Presentation


def contains_all(string, lst):
    """Returns True if and only if each meaningful (non-empty, longer than one char) element from the list-like lst
    is contained in the string"""
    for item in lst:
        item = item.strip('": ')
        if len(item) > 1 and item not in string:
            return False
    return True


def pull_images(pptx_file, image_dir,
                keystring=''  # all whitespace-separated fragments of this are to be contained in the image file name
                # as eligibility screen
                # if left empty, the value of title from pres properties is used automatically -- TOKENIZE?
                ):
    """Pulling relevant images picked based on :pptx_file: title or :keystring: from :image_dir: into the
    :pptx_file: pre-initialized with a title slide, one per slide"""

    pptx_file = pptx_file.replace('\\', '/')
    image_dir = image_dir.replace('\\', '/')
    print(f'Pulling from {image_dir}\ninto {pptx_file}')

    prs = Presentation(pptx_file)
    if not keystring:
        keystring = prs.core_properties.title
    if not keystring:
        keystring = input(
            '\nTitle is blank. Please enter the image screening keystring. \n'
            'Only image files containing all whitespace-separated fragments of the keystring in their names \n' 
            'will be pulled:\n')

    if len(prs.slides) > 1:
        proceed = input(
            'WARNING! More than one slide detected. Eligible images will be added to the end. Any symbol to proceed, '
            'mere Enter to cancel.'
        )
        if not proceed:
            print('Operation aborted. Please check destination file contents.')
            return

    blank_slide_layout = prs.slide_layouts[6]

    image_files = [
        file for file in os.scandir(image_dir) if contains_all(file.name, keystring.split(' ')
                                                               ) and file.name.endswith('.png')
    ]  # TODO: tokenize rather than split?
    image_files.sort(key=lambda x: x.stat().st_ctime)
    print(f'Total eligible images: {len(image_files)}')

    for image in image_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image.path, 0, 0, prs.slide_width)
        print('.', end="")
    prs.save(pptx_file)
    
    print('\nDone')


def main():
    pull_images(
        r'F:\User\Learn\ไทยศึกษา\เสียงสะท้อนอดีต\ของในหีบสำคัญอย่างไร สิ่งนี้ยืนยันความเชื่อเรื่องพระแก้วเมื่อ 500 ปีก่อน.pptx',
        r'C:\Users\User\Videos\Captures'
    )


if __name__ == '__main__':
    main()
