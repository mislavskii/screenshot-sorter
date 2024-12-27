import logging
import os
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format='%(levelname)s (%(module)s>%(name)s): %(message)s')


class Screens2PPTX:

    @staticmethod
    def contains_all(string, lst):
        """Returns True if and only if each meaningful (non-empty, longer than one char) element
        from the list-like lst is contained in the string"""
        for item in lst:
            item = item.strip('":. ')
            if len(item) > 1 and item not in string:
                return False
        return True

    def __init__(self, pptx_file, image_dir,
                 keystring=''  # all whitespace-separated fragments of this are to be contained in the image file name
                               # as eligibility screen. If left empty,
                               # the value of title from pres properties is used automatically - TOKENIZE?
                 ):
        self.pptx_file = pptx_file.replace('\\', '/')
        self.prs = Presentation(pptx_file)
        self.image_dir = image_dir.replace('\\', '/')
        self.keystring = keystring
        self.title = self.pptx_file.rsplit('/', 1)[-1].rsplit('.', 1)[0]
        self.prs.core_properties.title = self.title
        logging.info(f'initialized to pull images into {self.pptx_file} from {self.image_dir}')

    def populate_title(self):
        title_slide = self.prs.slides[0]
        title_shape = title_slide.shapes.title
        logging.info(f'setting title to "{self.title}"')
        title_shape.text = self.title

    def build_title_slide(self):
        logging.info('building the title slide')
        self.populate_title()
        logging.info('saving the presentation')
        self.prs.save(self.pptx_file)

    def pull_images(self):
        prs = self.prs
        keystring = self.keystring
        if len(prs.slides) > 1:
            proceed = input(
                'WARNING! More than one slide detected. Eligible images will be added to the end. Any symbol to '
                'proceed, mere Enter to cancel.'
            )
            if not proceed:
                print('Operation aborted. Please check destination file contents.')
                return
        blank_slide_layout = prs.slide_layouts[6]
        if not keystring:
            keystring = self.title
        logging.info(f'about to pull images with filtering on "{keystring}"')
        image_files = [file for file in os.scandir(self.image_dir) if self.contains_all(
            file.name, keystring.split(' ')) and file.name.endswith('.png')]  # TODO: tokenize rather than split?
        image_files.sort(key=lambda x: x.stat().st_ctime)
        logging.info(f'total eligible images: {len(image_files)}')
        for image in image_files:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(image.path, 0, 0, prs.slide_width)
            print('.', end="")
        print()
        logging.info('done pulling images.')
        prs.save(self.pptx_file)
        logging.info('presentation file saved.')
